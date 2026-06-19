# -*- coding: utf-8 -*-
"""Add real-world 2026 pricing slides to AriamindX_Hosting_Scaling.pptx.
Pricing sourced via web search, June 2026. See chat for source links."""
import copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SRC = "AriamindX_Hosting_Scaling.pptx"

# palette
BG      = RGBColor(0x0B, 0x0E, 0x17)
CARD    = RGBColor(0x15, 0x1B, 0x2B)
ROW_ALT = RGBColor(0x10, 0x14, 0x1F)
BLUE    = RGBColor(0x6E, 0x9B, 0xFF)
GREEN   = RGBColor(0x34, 0xD3, 0x99)
ORANGE  = RGBColor(0xFF, 0xB4, 0x4D)
RED     = RGBColor(0xFF, 0x6B, 0x6B)
WHITE   = RGBColor(0xF2, 0xF5, 0xFA)
MUTED   = RGBColor(0x9A, 0xA4, 0xB8)
DARK    = RGBColor(0x0B, 0x0E, 0x17)
FONT    = "Segoe UI"

prs = Presentation(SRC)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_rect(slide, x, y, w, h, fill, rounded=False, line=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Emu(int(x * 914400)), Emu(int(y * 914400)),
        Emu(int(w * 914400)), Emu(int(h * 914400)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=2):
    """runs: list of paragraphs; each paragraph is list of (text,size,bold,color)."""
    tb = slide.shapes.add_textbox(Emu(int(x * 914400)), Emu(int(y * 914400)),
                                  Emu(int(w * 914400)), Emu(int(h * 914400)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (text, size, bold, color) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = FONT
            r.font.color.rgb = color
    return tb


def base_slide(eyebrow, title):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, 13.333, 7.5, BG)                 # background
    add_rect(s, 0, 0, 0.18, 7.5, BLUE)                 # left accent bar
    add_text(s, 0.55, 0.35, 12.0, 0.4,
             [[(eyebrow, 13, True, BLUE)]])
    add_text(s, 0.5, 0.66, 12.3, 0.9,
             [[(title, 30, True, WHITE)]])
    return s


# ----------------------------------------------------------------------------
# SLIDE A : Service pricing table
# ----------------------------------------------------------------------------
sA = base_slide("WHAT IT ACTUALLY COSTS  —  2026 PRICING",
                "Price of the Recommended Stack")

tbl_x, tbl_w = 0.55, 12.23
c1 = tbl_x + 0.18          # service col
w1 = 3.65
c2 = c1 + w1               # plan col
w2 = 1.75
c3 = c2 + w2               # price col
w3 = tbl_x + tbl_w - c3 - 0.15

head_y = 1.7
row_h = 0.475
n_rows = 10

# header background
add_rect(sA, tbl_x, head_y, tbl_w, 0.42, BLUE, rounded=False)
add_text(sA, c1, head_y + 0.07, w1, 0.3, [[("SERVICE  /  ROLE", 11, True, DARK)]])
add_text(sA, c2, head_y + 0.07, w2, 0.3, [[("PLAN", 11, True, DARK)]])
add_text(sA, c3, head_y + 0.07, w3, 0.3, [[("PRICE  (June 2026)", 11, True, DARK)]])

rows = [
    ("Vercel", "App / web — serverless", "Pro",
     "$20", "/seat/mo · incl. $20 credit, 1 TB transfer, 10M edge req"),
    ("Google Cloud Run", "Containers / API tier", "Pay-go",
     "$0.40", "/M req + $24/M vCPU-s + $2.5/M GiB-s · 2M req/mo free"),
    ("Hetzner AX42", "Dedicated bare-metal box", "Monthly",
     "€39", "/mo · Ryzen 8700G, 64 GB DDR5, NVMe (+€39 setup)"),
    ("Neon", "Managed Postgres (pooled)", "Launch→Scale",
     "from $5", "/mo · $0.106→$0.222 / CU-hr + $0.35 / GB-mo"),
    ("Supabase", "Managed Postgres (alt)", "Pro",
     "$25", "/mo · 8 GB DB, 100 GB storage, 50K MAU incl."),
    ("Cloudflare R2", "Media / object storage", "Standard",
     "$0.015", "/GB-mo · $0 egress · $4.50/M Class A, $0.36/M Class B"),
    ("AWS S3", "Object storage (alt)", "Standard",
     "$0.023", "/GB-mo + $0.09/GB egress (first 100 GB/mo free)"),
    ("Upstash Redis", "Cache / rate-limit / queue", "Pay-go",
     "$0.20", "/100K commands · 256 MB + 500K cmd/mo free"),
    ("Cloudflare CDN", "Edge cache + WAF / DDoS", "Free / Pro",
     "$0 / $20", "/mo · unlimited cached bandwidth on free tier"),
    ("OpenAI / Anthropic /", "Gemini — AI inference", "Per-use",
     "pass-through", "· pay per token / image; the real variable cost"),
]

y = head_y + 0.42
for i, (name, role, plan, price, note) in enumerate(rows):
    fill = CARD if i % 2 == 0 else ROW_ALT
    add_rect(sA, tbl_x, y, tbl_w, row_h, fill)
    add_text(sA, c1, y + 0.05, w1, row_h,
             [[(name, 12.5, True, WHITE)], [(role, 10, False, MUTED)]],
             anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    add_text(sA, c2, y + 0.05, w2, row_h, [[(plan, 11, True, BLUE)]],
             anchor=MSO_ANCHOR.MIDDLE)
    pcolor = GREEN if i not in (9,) else ORANGE
    add_text(sA, c3, y + 0.05, w3, row_h,
             [[(price, 12.5, True, pcolor), (" ", 9, False, MUTED)],
              [(note.strip(), 10, False, MUTED)]],
             anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    y += row_h

add_text(sA, 0.55, 6.78, 12.3, 0.5,
         [[("Bottom line:  ", 12.5, True, GREEN),
           ("a 10K-user launch runs on managed services for roughly $45–$70/mo "
            "(or a €39 Hetzner box). AI API usage — not hosting — dominates the bill as you grow.",
            12.5, False, WHITE)]])

# ----------------------------------------------------------------------------
# SLIDE B : Cost by growth stage
# ----------------------------------------------------------------------------
sB = base_slide("BALLPARK MONTHLY BILL  —  HOSTING ONLY",
                "What You'll Pay at Each Stage")

cards = [
    ("10K", "users", "Launch", GREEN, "$45 – $70 / mo",
     ["Vercel Pro  $20", "Neon Launch  ~$20", "R2 media  $1–3",
      "Upstash + CDN  free", "— or — 1 Hetzner box €39"]),
    ("100K", "users", "Split & cache", BLUE, "$200 – $500 / mo",
     ["Cloud Run / box  $80+", "Neon Scale  ~$150", "R2 (100s GB)  $5–15",
      "Upstash paid  $20+", "Cloudflare Pro  $20"]),
    ("1M", "users", "Replicate", ORANGE, "$1.5K – $4K / mo",
     ["App fleet (Cloud Run)", "Neon Scale + replicas", "R2 (multi-TB) $30–100",
      "Redis cluster", "WAF + monitoring"]),
    ("10M", "users", "Fleet & harden", RED, "$10K – $30K+ / mo",
     ["Multi-box app fleet", "Large managed PG HA", "R2 (10s of TB)",
      "Dedicated workers", "Enterprise WAF / SLA"]),
]

cw, gap = 2.92, 0.18
cx0 = 0.55
cy = 1.75
ch = 4.55
for i, (num, unit, label, color, total, items) in enumerate(cards):
    x = cx0 + i * (cw + gap)
    add_rect(sB, x, cy, cw, ch, CARD, rounded=True)
    add_rect(sB, x, cy, cw, 0.16, color)                       # top accent strip
    add_text(sB, x + 0.22, cy + 0.32, cw - 0.4, 0.9,
             [[(num, 30, True, color), ("  " + unit, 13, True, MUTED)],
              [(label, 13, True, WHITE)]], space_after=2)
    add_text(sB, x + 0.22, cy + 1.45, cw - 0.4, 0.5,
             [[(total, 17, True, color)]])
    add_text(sB, x + 0.22, cy + 2.15, cw - 0.4, 2.2,
             [[("• " + it, 11.5, False, WHITE)] for it in items],
             space_after=5)

add_text(sB, 0.55, 6.62, 12.3, 0.7,
         [[("Order-of-magnitude estimates for infrastructure only.  ", 12, True, ORANGE),
           ("Excludes AI API spend, which scales with usage and typically exceeds hosting at every stage. "
            "Cloud prices change — verify before budgeting.", 12, False, MUTED)]])

# ----------------------------------------------------------------------------
# Move the two new slides to sit right before the final "Summary" slide
# ----------------------------------------------------------------------------
xml_slides = prs.slides._sldIdLst
ids = list(xml_slides)
newB = ids[-1]   # slide B (added last)
newA = ids[-2]   # slide A
summary = ids[-3]  # original last slide (Summary)
xml_slides.remove(newA)
xml_slides.remove(newB)
xml_slides.insert(list(xml_slides).index(summary), newA)
xml_slides.insert(list(xml_slides).index(summary), newB)

try:
    prs.save(SRC)
    print("Saved", SRC, "-> now", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
except PermissionError:
    alt = "AriamindX_Hosting_Scaling_priced.pptx"
    prs.save(alt)
    print("LOCKED: saved to", alt, "instead (close PowerPoint to overwrite original)")
