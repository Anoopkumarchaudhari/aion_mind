"""Generate the AriamindX hosting & scaling deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette ---------------------------------------------------------------
BG      = RGBColor(0x0B, 0x0E, 0x17)   # near-black navy
CARD    = RGBColor(0x15, 0x1B, 0x2B)
ACCENT  = RGBColor(0x6E, 0x8BFF & 0xFF, 0xFF)  # placeholder, set below
ACCENT  = RGBColor(0x6E, 0x9B, 0xFF)   # blue
ACCENT2 = RGBColor(0x34, 0xD3, 0x99)   # green
WARN    = RGBColor(0xFF, 0xB4, 0x4D)   # amber
DANGER  = RGBColor(0xFF, 0x6B, 0x6B)   # red
WHITE   = RGBColor(0xF2, 0xF5, 0xFA)
MUTED   = RGBColor(0x9A, 0xA4, 0xB8)

WIDE_W, WIDE_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = WIDE_W
prs.slide_height = WIDE_H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, WIDE_W, WIDE_H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=True):
    shp = s.shapes.add_shape(5 if radius else 1, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    if isinstance(runs, str):
        runs = [[(runs, 18, WHITE, False)]]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        if isinstance(para, tuple):
            para = [para]
        for seg in para:
            t, sz, col, bold = (list(seg) + [False])[:4]
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = "Segoe UI"
    return tb


def header(s, kicker, title):
    box(s, Inches(0), Inches(0), Inches(0.18), WIDE_H, fill=ACCENT, radius=False)
    text(s, Inches(0.55), Inches(0.35), Inches(12), Inches(0.4),
         [[(kicker, 13, ACCENT, True)]])
    text(s, Inches(0.5), Inches(0.66), Inches(12.3), Inches(0.9),
         [[(title, 30, WHITE, True)]])


def chip(s, x, y, label, col):
    w = Inches(0.13 + 0.092 * len(label))
    c = box(s, x, y, w, Inches(0.34), fill=CARD, line=col, line_w=1.25)
    text(s, x, y, w, Inches(0.34), [[(label, 11, col, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    return x + w + Inches(0.12)


# ===========================================================================
# 1. TITLE
# ===========================================================================
s = slide()
box(s, Inches(0), Inches(2.55), WIDE_W, Inches(0.045), fill=ACCENT, radius=False)
text(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.0),
     [[("AriamindX", 54, WHITE, True)]])
text(s, Inches(0.92), Inches(2.75), Inches(11.5), Inches(0.7),
     [[("Hosting, Storage & Scaling Architecture", 26, ACCENT, True)]])
text(s, Inches(0.92), Inches(3.5), Inches(11.5), Inches(0.6),
     [[("Scaling plan for 10K → 10M users  •  Next.js 15 + PostgreSQL + external AI APIs", 16, MUTED, False)]])
xx = Inches(0.92); y = Inches(4.7)
for lab, c in [("Next.js 15", ACCENT), ("PostgreSQL", ACCENT2),
               ("Object Storage", WARN), ("CDN", ACCENT), ("No GPU needed", ACCENT2)]:
    xx = chip(s, xx, y, lab, c)
text(s, Inches(0.92), Inches(6.6), Inches(11), Inches(0.4),
     [[("Prepared for self-managed deployment planning  •  2026", 12, MUTED, False)]])

# ===========================================================================
# 2. BEST SETUP (RECOMMENDED) -- on top
# ===========================================================================
s = slide()
header(s, "START HERE — THE RECOMMENDATION", "Best Setup (Recommended)")

# big recommended bar
box(s, Inches(0.55), Inches(1.65), Inches(12.25), Inches(1.1), fill=CARD, line=ACCENT2, line_w=2)
text(s, Inches(0.8), Inches(1.78), Inches(11.8), Inches(0.9),
     [[("Front end on managed platform  +  Pooled Postgres  +  Object storage + CDN", 20, WHITE, True)],
      [("Vercel / Cloud Run  •  Neon (pooled)  •  Cloudflare R2  •  Upstash Redis  •  Cloudflare CDN", 14, ACCENT2, True)]],
     space_after=4)

cards = [
    ("App / Web tier", ACCENT, ["Vercel (now) → Cloud Run", "or a Hetzner dedicated box", "Auto-TLS, scales horizontally"]),
    ("Database", ACCENT2, ["Neon / Supabase Postgres", "Pooled connection string", "Chat + credentials + credits"]),
    ("Media files", WARN, ["Cloudflare R2 (zero egress)", "Generated images & videos", "Served via CDN, not API"]),
    ("Cache / Queue", ACCENT, ["Upstash Redis", "Rate-limit paid AI calls", "Background media jobs"]),
]
cw = Inches(2.92); gap = Inches(0.16); x0 = Inches(0.55); cy = Inches(3.0); ch = Inches(2.65)
for i, (t, col, items) in enumerate(cards):
    x = x0 + i * (cw + gap)
    box(s, x, cy, cw, ch, fill=CARD, line=col, line_w=1.5)
    box(s, x, cy, cw, Inches(0.12), fill=col, radius=False)
    text(s, x + Inches(0.18), cy + Inches(0.22), cw - Inches(0.3), Inches(0.5),
         [[(t, 16, col, True)]])
    text(s, x + Inches(0.18), cy + Inches(0.78), cw - Inches(0.3), Inches(1.7),
         [[("• " + it, 12.5, WHITE, False)] for it in items], space_after=7, line_spacing=1.0)

text(s, Inches(0.55), Inches(5.95), Inches(12.2), Inches(1.2),
     [[("Why this wins:  ", 14, ACCENT2, True),
       ("cheapest path that survives every growth stage. Compute scales horizontally, the DB is "
        "managed + pooled, and big files never touch your own bandwidth. Same stack works at 10K and at 10M — "
        "you only change tiers, not architecture.", 14, WHITE, False)]],
     line_spacing=1.05)

# ===========================================================================
# 3. CURRENT ARCHITECTURE
# ===========================================================================
s = slide()
header(s, "WHERE YOU ARE TODAY", "Current Architecture")
rows = [
    ("Framework", "Next.js 15 (App Router) + React 19", "Scales well, serverless-friendly", ACCENT2, "OK"),
    ("Database", "Raw PostgreSQL via pg.Pool  (services/db.ts)", "Pool-per-instance + DDL on first query", WARN, "WATCH"),
    ("Auth", "DB session table + Google OAuth (services/auth.ts)", "Stateful in Postgres — fine", ACCENT2, "OK"),
    ("File storage", "Local disk — data/  (generatedImageFiles.ts)", "Breaks on multi-instance / redeploy", DANGER, "FIX"),
    ("AI / media", "OpenAI, Anthropic, Gemini, DeepSeek, Runware", "External APIs — cost & latency live here", ACCENT2, "OK"),
    ("Payments/email", "Razorpay + nodemailer", "Fine", ACCENT2, "OK"),
]
y = Inches(1.7); rh = Inches(0.82)
# header row
text(s, Inches(0.6), y, Inches(2.4), rh, [[("LAYER", 12, MUTED, True)]])
text(s, Inches(3.0), y, Inches(4.6), rh, [[("WHAT YOU USE", 12, MUTED, True)]])
text(s, Inches(7.7), y, Inches(4.2), rh, [[("SCALING IMPLICATION", 12, MUTED, True)]])
text(s, Inches(11.9), y, Inches(1.3), rh, [[("STATUS", 12, MUTED, True)]])
y = y + Inches(0.42)
for layer, use, impl, col, status in rows:
    box(s, Inches(0.55), y, Inches(12.25), rh - Inches(0.08), fill=CARD)
    box(s, Inches(0.55), y, Inches(0.08), rh - Inches(0.08), fill=col, radius=False)
    text(s, Inches(0.75), y, Inches(2.3), rh - Inches(0.08), [[(layer, 13, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.0), y, Inches(4.6), rh - Inches(0.08), [[(use, 12, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(7.7), y, Inches(4.1), rh - Inches(0.08), [[(impl, 12, WHITE, False)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(11.85), y, Inches(0.95), rh - Inches(0.08), [[(status, 12, col, True)]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    y = y + rh

# ===========================================================================
# 4. BLOCKERS IN CODE
# ===========================================================================
s = slide()
header(s, "MUST FIX BEFORE SCALING", "Three Blockers in Your Code")
blk = [
    ("1", "Local filesystem storage", DANGER,
     "generatedImageFiles.ts writes images to process.cwd()/data/. On serverless or "
     "multi-instance hosts this is ephemeral and not shared — files vanish on redeploy and are "
     "invisible to other instances.",
     "Move to Cloudflare R2 / S3 / GCS."),
    ("2", "pg.Pool per instance", WARN,
     "Each instance opens its own pool (db.ts). 100 instances × pool size = thousands of DB "
     "connections → Postgres falls over.",
     "Connect through a pooler (Neon / PgBouncer / Supabase)."),
    ("3", "DDL on every cold start", WARN,
     "ensureDatabaseSchema() runs CREATE TABLE on first query. Fine at 10K; a race + latency hit at scale.",
     "Run migrations once at deploy time."),
]
y = Inches(1.75)
for n, t, col, body, fix in blk:
    box(s, Inches(0.55), y, Inches(12.25), Inches(1.55), fill=CARD, line=col, line_w=1.4)
    box(s, Inches(0.55), y, Inches(0.9), Inches(1.55), fill=col, radius=False)
    text(s, Inches(0.55), y, Inches(0.9), Inches(1.55), [[(n, 40, BG, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.65), y + Inches(0.12), Inches(11), Inches(0.5), [[(t, 18, col, True)]])
    text(s, Inches(1.65), y + Inches(0.58), Inches(10.9), Inches(0.7), [[(body, 12.5, WHITE, False)]], line_spacing=1.0)
    text(s, Inches(1.65), y + Inches(1.16), Inches(10.9), Inches(0.35),
         [[("Fix:  ", 12.5, ACCENT2, True), (fix, 12.5, ACCENT2, False)]])
    y = y + Inches(1.72)

# ===========================================================================
# 5. HOSTING OPTIONS COMPARISON
# ===========================================================================
s = slide()
header(s, "PICK YOUR COMPUTE", "Hosting Options Compared")
cols = [
    ("Vercel / Netlify", ACCENT, "Managed serverless",
     ["Zero-config, auto-scale", "Best for 10K–500K", "Watch function timeouts", "$$ at high volume"]),
    ("Cloud Run / Fargate", ACCENT2, "Managed containers",
     ["Autoscale + long requests", "Sweet spot 100K–10M", "No cold-start tax", "Best for media jobs"]),
    ("VPS / Dedicated", WARN, "Self-managed",
     ["Cheapest (Hetzner ~€50)", "Full control", "You add LB + replicas", "Manual scaling"]),
    ("Home machine", DANGER, "Own hardware",
     ["Dev / ≤ ~1–2K users only", "CGNAT, dynamic IP", "No redundancy / SLA", "DDoS exposure"]),
]
cw = Inches(2.92); gap = Inches(0.16); x0 = Inches(0.55); cy = Inches(1.75); ch = Inches(4.6)
for i, (t, col, sub, items) in enumerate(cols):
    x = x0 + i * (cw + gap)
    box(s, x, cy, cw, ch, fill=CARD, line=col, line_w=1.5)
    box(s, x, cy, cw, Inches(0.7), fill=col, radius=True)
    box(s, x, cy + Inches(0.35), cw, Inches(0.35), fill=col, radius=False)
    text(s, x, cy + Inches(0.06), cw, Inches(0.4), [[(t, 15, BG, True)]], align=PP_ALIGN.CENTER)
    text(s, x, cy + Inches(0.4), cw, Inches(0.3), [[(sub, 11, BG, True)]], align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.18), cy + Inches(0.92), cw - Inches(0.32), Inches(3.5),
         [[("• " + it, 12.5, WHITE, False)] for it in items], space_after=10, line_spacing=1.0)
text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.6),
     [[("Recommended path:  ", 13, ACCENT2, True),
       ("Start on Vercel → migrate the API + worker tier to Cloud Run (or a Hetzner dedicated box) around 100K–500K users.", 13, WHITE, False)]])

# ===========================================================================
# 6. STORAGE STRATEGY
# ===========================================================================
s = slide()
header(s, "RIGHT DATA, RIGHT STORE", "Storage Strategy")
# two big columns
box(s, Inches(0.55), Inches(1.75), Inches(6.0), Inches(4.7), fill=CARD, line=ACCENT2, line_w=1.5)
box(s, Inches(0.55), Inches(1.75), Inches(6.0), Inches(0.6), fill=ACCENT2, radius=True)
box(s, Inches(0.55), Inches(2.1), Inches(6.0), Inches(0.25), fill=ACCENT2, radius=False)
text(s, Inches(0.55), Inches(1.83), Inches(6.0), Inches(0.45), [[("PostgreSQL  —  structured data", 16, BG, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.8), Inches(2.6), Inches(5.5), Inches(3.6),
     [[("✓ Chat threads & messages", 14, WHITE, True)],
      [("✓ User credentials & password hashes", 14, WHITE, True)],
      [("✓ Sessions, credits, payments ledger", 14, WHITE, True)],
      [("", 6, WHITE, False)],
      [("Why: transactions, indexes, foreign keys,", 12.5, MUTED, False)],
      [("UNIQUE(email), row locking, instant lookups.", 12.5, MUTED, False)],
      [("", 6, WHITE, False)],
      [("Managed: Neon / Supabase / Cloud SQL", 13, ACCENT2, True)],
      [("Always via a connection pooler.", 12.5, ACCENT2, False)]],
     space_after=6, line_spacing=1.0)

box(s, Inches(6.8), Inches(1.75), Inches(6.0), Inches(4.7), fill=CARD, line=WARN, line_w=1.5)
box(s, Inches(6.8), Inches(1.75), Inches(6.0), Inches(0.6), fill=WARN, radius=True)
box(s, Inches(6.8), Inches(2.1), Inches(6.0), Inches(0.25), fill=WARN, radius=False)
text(s, Inches(6.8), Inches(1.83), Inches(6.0), Inches(0.45), [[("Object storage  —  big files", 16, BG, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(7.05), Inches(2.6), Inches(5.5), Inches(3.6),
     [[("✓ Generated images", 14, WHITE, True)],
      [("✓ Generated videos", 14, WHITE, True)],
      [("✓ Uploaded assets / sidebar media", 14, WHITE, True)],
      [("", 6, WHITE, False)],
      [("Why: cheap per-GB, infinite scale, served", 12.5, MUTED, False)],
      [("by CDN — never proxied through your API.", 12.5, MUTED, False)],
      [("", 6, WHITE, False)],
      [("Best: Cloudflare R2 (zero egress fees)", 13, WARN, True)],
      [("Alternatives: AWS S3 / Google GCS", 12.5, WARN, False)]],
     space_after=6, line_spacing=1.0)

text(s, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.5),
     [[("Golden rule:  ", 14, DANGER, True),
       ("queried / transactional data → Postgres.   Large blobs → object storage.   Never the other way around.", 14, WHITE, False)]])

# ===========================================================================
# 7. DO YOU NEED A GPU
# ===========================================================================
s = slide()
header(s, "COMMON QUESTION", "Do You Need a GPU?  —  No")
box(s, Inches(0.55), Inches(1.75), Inches(12.25), Inches(1.0), fill=CARD, line=ACCENT2, line_w=2)
text(s, Inches(0.8), Inches(1.9), Inches(11.8), Inches(0.8),
     [[("All AI runs on external APIs. Your server only runs Next.js + PostgreSQL — ", 16, WHITE, False),
       ("CPU / RAM / disk / network bound, not GPU bound.", 16, ACCENT2, True)]], line_spacing=1.05)
rows = [
    ("Chat / image / video generation", "External APIs (pay per token/image)", "No", ACCENT2),
    ("Next.js app + PostgreSQL", "Your machine", "No", ACCENT2),
    ("Self-hosting your own LLM (e.g. Llama)", "Would need 2× A100/H100 — separate project", "Yes", DANGER),
]
y = Inches(3.1); rh = Inches(0.9)
text(s, Inches(0.75), y, Inches(5.5), Inches(0.4), [[("WORKLOAD", 12, MUTED, True)]])
text(s, Inches(6.3), y, Inches(5.0), Inches(0.4), [[("WHERE IT RUNS", 12, MUTED, True)]])
text(s, Inches(11.4), y, Inches(1.4), Inches(0.4), [[("GPU?", 12, MUTED, True)]])
y = y + Inches(0.45)
for w, where, gpu, col in rows:
    box(s, Inches(0.55), y, Inches(12.25), rh - Inches(0.1), fill=CARD)
    box(s, Inches(0.55), y, Inches(0.08), rh - Inches(0.1), fill=col, radius=False)
    text(s, Inches(0.75), y, Inches(5.5), rh - Inches(0.1), [[(w, 13.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(6.3), y, Inches(5.0), rh - Inches(0.1), [[(where, 12.5, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(11.3), y, Inches(1.4), rh - Inches(0.1), [[(gpu, 14, col, True)]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    y = y + rh
text(s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.5),
     [[("Spend the GPU budget on RAM, NVMe SSD and bandwidth instead.", 14, ACCENT, True)]])

# ===========================================================================
# 8. SELF-HOST HARDWARE REQUIREMENTS
# ===========================================================================
s = slide()
header(s, "IF YOU MANAGE YOUR OWN SYSTEM", "Hardware Requirements by Scale")
tiers = [
    ("Dev / ≤1K", ACCENT, ["CPU: 4 cores", "RAM: 8–16 GB", "Disk: 256 GB SSD", "Net: home broadband", "DB: same box"]),
    ("~10K", ACCENT, ["CPU: 4–8 cores", "RAM: 16–32 GB", "Disk: 500 GB NVMe", "Net: 1 Gbps + static IP", "DB: same box OK"]),
    ("~100K", ACCENT2, ["App: 8c / 32 GB", "DB: 8–16c / 64 GB", "NVMe high IOPS", "+ Redis + LB", "App & DB split"]),
    ("1M–10M", WARN, ["4–8 app boxes", "DB: 32–64c / 128–256 GB", "Primary + replicas", "Redis cluster + workers", "10 Gbps + DDoS"]),
]
cw = Inches(2.92); gap = Inches(0.16); x0 = Inches(0.55); cy = Inches(1.8); ch = Inches(3.7)
for i, (t, col, items) in enumerate(tiers):
    x = x0 + i * (cw + gap)
    box(s, x, cy, cw, ch, fill=CARD, line=col, line_w=1.5)
    box(s, x, cy, cw, Inches(0.6), fill=col, radius=True)
    box(s, x, cy + Inches(0.3), cw, Inches(0.3), fill=col, radius=False)
    text(s, x, cy + Inches(0.08), cw, Inches(0.45), [[(t + " users", 15, BG, True)]], align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.2), cy + Inches(0.78), cw - Inches(0.35), Inches(2.8),
         [[("• " + it, 12.5, WHITE, False)] for it in items], space_after=9, line_spacing=1.0)
text(s, Inches(0.55), Inches(5.75), Inches(12.2), Inches(1.4),
     [[("Practical middle ground:  ", 14, ACCENT2, True),
       ("rent a dedicated bare-metal server (Hetzner ~€40–60/mo: Ryzen, 64 GB RAM, 2×NVMe, 1 Gbps "
        "unmetered). Same “my own system” control, real bandwidth + static IP, no home-hosting blockers. "
        "One such box handles tens of thousands of users.", 13.5, WHITE, False)]],
     line_spacing=1.1)

# ===========================================================================
# 9. BANDWIDTH & NETWORK
# ===========================================================================
s = slide()
header(s, "THE REAL CONSTRAINT", "Bandwidth & Network")
box(s, Inches(0.55), Inches(1.75), Inches(6.0), Inches(4.6), fill=CARD, line=ACCENT, line_w=1.5)
text(s, Inches(0.8), Inches(1.9), Inches(5.6), Inches(0.5), [[("The math", 16, ACCENT, True)]])
text(s, Inches(0.8), Inches(2.45), Inches(5.5), Inches(3.6),
     [[("• Chat JSON: ~50 KB / response", 13.5, WHITE, False)],
      [("• 10K concurrent ≈ hundreds of Mbps bursts → needs 1 Gbps", 13.5, WHITE, False)],
      [("• Serving a 5 MB image to 1,000 users = 5 GB of egress", 13.5, WHITE, False)],
      [("• Media from your own box = 50–200× the traffic", 13.5, DANGER, True)],
      [("", 8, WHITE, False)],
      [("Fix: put media on R2 + CDN and let them eat the egress.", 13.5, ACCENT2, True)]],
     space_after=9, line_spacing=1.05)

box(s, Inches(6.8), Inches(1.75), Inches(6.0), Inches(4.6), fill=CARD, line=WARN, line_w=1.5)
text(s, Inches(7.05), Inches(1.9), Inches(5.6), Inches(0.5), [[("Network checklist", 16, WARN, True)]])
text(s, Inches(7.05), Inches(2.45), Inches(5.5), Inches(3.6),
     [[("✓ Static public IP (not CGNAT)", 13.5, WHITE, False)],
      [("✓ Ports 80 / 443 open & ISP-allowed", 13.5, WHITE, False)],
      [("✓ High / symmetric upload (business fiber)", 13.5, WHITE, False)],
      [("✓ UPS for power; ideally 2nd ISP for failover", 13.5, WHITE, False)],
      [("✓ Cloudflare in front: TLS, cache, DDoS shield", 13.5, WHITE, False)]],
     space_after=11, line_spacing=1.05)
text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.6),
     [[("Residential upload (10–40 Mbps) + CGNAT + no redundancy = why home hosting cannot reach real scale.", 13.5, DANGER, True)]])

# ===========================================================================
# 10. STEP-BY-STEP SCALING ROADMAP
# ===========================================================================
s = slide()
header(s, "THE GROWTH PATH", "Step-by-Step Scaling Roadmap")
steps = [
    ("10K", ACCENT, "Launch simple",
     "Vercel (or 1 Hetzner box) + managed Neon Postgres (pooled) + R2 for media + Cloudflare CDN. App + DB can share resources."),
    ("100K", ACCENT, "Split & cache",
     "Separate app and DB. Add Redis (cache + rate-limit). Move media generation to an async job queue. CDN on all static + media."),
    ("1M", ACCENT2, "Replicate & isolate",
     "Postgres primary + read replicas (chat reads → replicas). Split web tier from media-worker tier. Autoscale containers on Cloud Run / Fargate."),
    ("10M", WARN, "Fleet & harden",
     "Multi-instance app fleet behind LB, Redis cluster, dedicated workers, 10 Gbps, Cloudflare WAF + DDoS. Effectively a small datacenter — rent, don't build."),
]
y = Inches(1.8); rh = Inches(1.18)
for i, (tag, col, title, body) in enumerate(steps):
    box(s, Inches(0.55), y, Inches(12.25), rh - Inches(0.12), fill=CARD, line=col, line_w=1.3)
    box(s, Inches(0.55), y, Inches(1.7), rh - Inches(0.12), fill=col, radius=True)
    box(s, Inches(1.5), y, Inches(0.75), rh - Inches(0.12), fill=col, radius=False)
    text(s, Inches(0.55), y, Inches(1.7), rh - Inches(0.12), [[(tag, 26, BG, True)], [("users", 11, BG, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, Inches(2.45), y + Inches(0.13), Inches(10.1), Inches(0.4), [[(title, 16, col, True)]])
    text(s, Inches(2.45), y + Inches(0.52), Inches(10.2), Inches(0.6), [[(body, 12.5, WHITE, False)]], line_spacing=1.0)
    if i < len(steps) - 1:
        text(s, Inches(1.32), y + rh - Inches(0.2), Inches(0.4), Inches(0.3), [[("▼", 12, col, True)]], align=PP_ALIGN.CENTER, space_after=0)
    y = y + rh

# ===========================================================================
# 11. MIGRATION CHECKLIST (priority)
# ===========================================================================
s = slide()
header(s, "DO THESE IN ORDER", "Migration Checklist")
items = [
    (DANGER, "1", "Move generated image/video storage from local fs → Cloudflare R2 / S3", "Unblocks horizontal scaling"),
    (DANGER, "2", "Move Postgres to Neon / Supabase, connect via pooled URL", "Fixes connection exhaustion"),
    (WARN,   "3", "Replace DDL-on-startup with a migrations step at deploy", "Removes cold-start race"),
    (WARN,   "4", "Make image/video generation async via a queue (Redis)", "Long jobs stop blocking requests"),
    (WARN,   "5", "Add Cloudflare CDN + rate-limiting on paid AI endpoints", "Protects cost & latency"),
    (ACCENT, "6", "At ~500K: split web vs media-worker tiers, add read replicas", "Independent scaling"),
]
y = Inches(1.8); rh = Inches(0.8)
for col, n, t, why in items:
    box(s, Inches(0.55), y, Inches(12.25), rh - Inches(0.12), fill=CARD)
    box(s, Inches(0.55), y, Inches(0.7), rh - Inches(0.12), fill=col, radius=False)
    text(s, Inches(0.55), y, Inches(0.7), rh - Inches(0.12), [[(n, 22, BG, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.45), y, Inches(8.3), rh - Inches(0.12), [[(t, 14, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(9.9), y, Inches(2.8), rh - Inches(0.12), [[(why, 11.5, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
    y = y + rh

# ===========================================================================
# 12. SUMMARY
# ===========================================================================
s = slide()
header(s, "ONE PAGE TO REMEMBER", "Summary")
pts = [
    (ACCENT2, "Best stack", "Vercel/Cloud Run + Neon (pooled) + Cloudflare R2 + Upstash Redis + Cloudflare CDN. Same architecture from 10K to 10M."),
    (DANGER,  "No GPU", "All AI is external APIs. Invest in RAM, NVMe and bandwidth, not GPUs."),
    (WARN,    "Storage rule", "Chat + credentials → PostgreSQL. Images + videos → object storage. Never swap them."),
    (ACCENT,  "Own system", "Home = dev / small only. For self-managed scale, rent a Hetzner-class dedicated box, not a home PC."),
    (ACCENT2, "First move", "Refactor local-disk media storage to pluggable R2/S3 — the single change that makes the app multi-instance ready."),
]
# (brand: AriamindX)
y = Inches(1.85); rh = Inches(0.95)
for col, t, body in pts:
    box(s, Inches(0.55), y, Inches(12.25), rh - Inches(0.13), fill=CARD, line=col, line_w=1.2)
    box(s, Inches(0.55), y, Inches(0.12), rh - Inches(0.13), fill=col, radius=False)
    text(s, Inches(0.85), y, Inches(2.6), rh - Inches(0.13), [[(t, 15, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.5), y, Inches(9.1), rh - Inches(0.13), [[(body, 13, WHITE, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    y = y + rh
text(s, Inches(0.55), Inches(6.95), Inches(12), Inches(0.4),
     [[("AriamindX  •  Hosting & Scaling Architecture", 11, MUTED, False)]])

import os
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "AriamindX_Hosting_Scaling.pptx")
prs.save(out)
print("Saved:", out, "slides:", len(prs.slides._sldIdLst))
